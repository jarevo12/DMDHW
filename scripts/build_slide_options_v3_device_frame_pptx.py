from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0, 0, 0)
BG_SECONDARY = RGBColor(10, 10, 10)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(192, 192, 192)
TEXT_MUTED = RGBColor(102, 102, 102)
ACCENT_PURPLE = RGBColor(167, 139, 250)
ACCENT_GREEN = RGBColor(204, 255, 0)

# Fonts
FONT_SANS = "Arial Black"
FONT_BODY = "Arial"
FONT_MONO = "Consolas"

# Layout constants
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
CONTENT_W = SLIDE_W - (MARGIN_X * 2)
GAP_COL = Inches(0.6)
COL_W = (CONTENT_W - GAP_COL) / 2


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, uppercase=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=None, line_width=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.color.rgb = border
    shape.line.width = Pt(line_width)
    shape.fill.solid()
    if fill_color is None:
        shape.fill.fore_color.rgb = BG
    else:
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_rounded_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=None, line_width=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.line.color.rgb = border
    shape.line.width = Pt(line_width)
    shape.fill.solid()
    if fill_color is None:
        shape.fill.fore_color.rgb = BG
    else:
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_logo(slide, x, y, size):
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    base.fill.solid()
    base.fill.fore_color.rgb = TEXT_PRIMARY
    base.line.color.rgb = TEXT_PRIMARY
    base.line.width = Pt(1)
    bar_h = size * 0.16
    bar_w1 = size * 0.7
    bar_w2 = size * 0.45
    bar_w3 = size * 0.18
    offset_x = x + size * 0.15
    base_y = y + size * 0.62
    mid_y = y + size * 0.42
    top_y = y + size * 0.22
    bars = [
        (offset_x, base_y, bar_w1, bar_h),
        (x + size * 0.28, mid_y, bar_w2, bar_h),
        (x + size * 0.41, top_y, bar_w3, bar_h),
    ]
    for bx, by, bw, bh in bars:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bw, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = BG
        bar.line.fill.background()


def add_section_label(slide, text, x, y):
    return add_textbox(slide, x, y, Inches(4), Inches(0.3), f"// {text}", font_name=FONT_MONO, font_size=14, bold=True, color=ACCENT_PURPLE)


def add_h1(slide, text, x, y, w):
    return add_textbox(slide, x, y, w, Inches(0.9), text, font_name=FONT_SANS, font_size=42, bold=True, color=TEXT_PRIMARY)


def add_h3(slide, text, x, y, w, color=TEXT_PRIMARY):
    return add_textbox(slide, x, y, w, Inches(0.4), text, font_name=FONT_SANS, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_device_phone(slide, x, y, frame_w, frame_h, border_color, play_color, label_text):
    # Shadow
    shadow = add_rounded_box(slide, x + Inches(0.08), y + Inches(0.1), frame_w, frame_h, border=BG, fill_color=BG, line_width=0.5)
    shadow.fill.transparency = 0.35
    shadow.line.fill.background()
    # Device frame (gradient approximated with dark fill)
    frame = add_rounded_box(slide, x, y, frame_w, frame_h, border=BG, fill_color=RGBColor(26, 26, 26), line_width=1)
    frame.line.fill.background()
    # Phone mockup
    padding = Inches(0.12)
    phone_x = x + padding
    phone_y = y + padding
    phone_w = frame_w - padding * 2
    phone_h = frame_h - padding * 2
    phone = add_rounded_box(slide, phone_x, phone_y, phone_w, phone_h, border=border_color, fill_color=BG_SECONDARY, line_width=2.5)
    # Notch
    notch = add_rounded_box(slide, phone_x + phone_w * 0.32, phone_y + phone_h * 0.03, phone_w * 0.36, phone_h * 0.06, border=BG, fill_color=BG, line_width=0.5)
    notch.line.fill.background()
    # Play button
    play_size = phone_w * 0.36
    play_x = phone_x + (phone_w - play_size) / 2
    play_y = phone_y + (phone_h - play_size) / 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, play_x, play_y, play_size, play_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_SECONDARY
    circle.line.color.rgb = play_color
    circle.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, play_x + play_size * 0.38, play_y + play_size * 0.28, play_size * 0.28, play_size * 0.32)
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = play_color
    tri.line.fill.background()
    # Label
    add_textbox(slide, phone_x, phone_y + phone_h - Inches(0.45), phone_w, Inches(0.3), label_text, font_name=FONT_BODY, font_size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_textbox(slide, MARGIN_X, Inches(0.35), Inches(7.2), Inches(0.3), "DEMO SLIDE - OPTION B: DEVICE FRAMES", font_name=FONT_BODY, font_size=11, bold=True, color=BG, align=PP_ALIGN.LEFT)
label_bg = add_box(slide, MARGIN_X, Inches(0.35), Inches(7.2), Inches(0.3), border=ACCENT_PURPLE, fill_color=ACCENT_PURPLE, line_width=1)
label_tf = label_bg.text_frame
label_tf.clear()
label_run = label_tf.paragraphs[0].add_run()
label_run.text = "DEMO SLIDE - OPTION B: DEVICE FRAMES"
label_run.font.name = FONT_BODY
label_run.font.size = Pt(11)
label_run.font.bold = True
label_run.font.color.rgb = BG
label_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

add_section_label(slide, "Product Demo", MARGIN_X, Inches(0.8))
add_h1(slide, "Experience Axiom Forge in Action", MARGIN_X, Inches(1.3), CONTENT_W)

card_y = Inches(2.35)
frame_w = Inches(3.0)
frame_h = Inches(5.0)
for i, (title, border_color, play_color, label_text) in enumerate([
    ("Onboarding Flow", ACCENT_GREEN, ACCENT_GREEN, "Watch: 45 sec"),
    ("App Features", ACCENT_PURPLE, ACCENT_PURPLE, "Watch: 2 min"),
]):
    col_x = MARGIN_X + i * (COL_W + GAP_COL)
    add_h3(slide, title, col_x, card_y, COL_W, color=ACCENT_GREEN if i == 0 else ACCENT_PURPLE)
    frame_x = col_x + (COL_W - frame_w) / 2
    frame_y = card_y + Inches(0.45)
    add_device_phone(slide, frame_x, frame_y, frame_w, frame_h, border_color, play_color, label_text)

# CTA box
cta_w = Inches(6.4)
cta_h = Inches(0.95)
cta_x = (SLIDE_W - cta_w) / 2
cta_y = Inches(6.2)
cta = add_box(slide, cta_x, cta_y, cta_w, cta_h, border=ACCENT_GREEN, fill_color=RGBColor(20, 26, 0), line_width=2)
add_textbox(slide, cta_x, cta_y + Inches(0.1), cta_w, Inches(0.3), "Try It Yourself", font_name=FONT_SANS, font_size=16, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_textbox(slide, cta_x, cta_y + Inches(0.45), cta_w, Inches(0.25), "No download required - works in your browser", font_name=FONT_BODY, font_size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
button = add_box(slide, cta_x + Inches(2.2), cta_y + Inches(0.62), Inches(2.0), Inches(0.3), border=ACCENT_GREEN, fill_color=ACCENT_GREEN, line_width=1)
btn_tf = button.text_frame
btn_tf.clear()
btn_run = btn_tf.paragraphs[0].add_run()
btn_run.text = "LAUNCH APP"
btn_run.font.name = FONT_SANS
btn_run.font.size = Pt(12)
btn_run.font.bold = True
btn_run.font.color.rgb = BG
btn_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

prs.save("/workspaces/DMDHW/MIT_Sandbox/Pitch Deck/slide-options-v3-device-frames.pptx")
print("Saved device frame slide PPTX")
