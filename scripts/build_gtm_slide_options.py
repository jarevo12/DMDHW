from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0, 0, 0)
BG_SECONDARY = RGBColor(10, 10, 10)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(192, 192, 192)
TEXT_MUTED = RGBColor(102, 102, 102)
ACCENT_PURPLE = RGBColor(167, 139, 250)
ACCENT_GREEN = RGBColor(204, 255, 0)
ACCENT_RED = RGBColor(255, 0, 0)
ACCENT_YELLOW = RGBColor(255, 255, 0)

# Fonts
FONT_SANS = "Arial Black"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"

# Layout constants
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
CONTENT_W = SLIDE_W - (MARGIN_X * 2)

def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box

def add_box(slide, x, y, w, h, border=None, fill_color=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
        
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_section_label(slide, text, x, y):
    add_textbox(slide, x, y, Inches(6), Inches(0.3), f"// {text.upper()}", font_name=FONT_MONO, font_size=14, bold=True, color=ACCENT_PURPLE)

def add_logo(slide):
    # Simplified logo representation (top right)
    size = Inches(0.5)
    x = SLIDE_W - MARGIN_X - size
    y = Inches(0.4)
    
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    base.fill.solid()
    base.fill.fore_color.rgb = TEXT_PRIMARY
    
    # Bars
    bar_h = size * 0.16
    bar_w1 = size * 0.7
    bar_w2 = size * 0.45
    bar_w3 = size * 0.18
    offset_x = x + size * 0.15
    base_y = y + size * 0.62
    mid_y = y + size * 0.42
    top_y = y + size * 0.22
    
    for bx, by, bw in [(offset_x, base_y, bar_w1), (x + size * 0.28, mid_y, bar_w2), (x + size * 0.41, top_y, bar_w3)]:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bw, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = BG
        bar.line.fill.background()

def add_bullets(slide, x, y, w, h, items, font_size=14, color=TEXT_SECONDARY, bullet_color=ACCENT_GREEN):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.space_before = Pt(6)
        # Custom bullet not easily doable with high fidelity in simple mode, stick to default or manual
        # We can simulate the square bullet style from CSS if we want, but standard bullet is fine for now.

# ==============================================================================
# SLIDE 1: INTRO
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, MARGIN_X, Inches(2.5), CONTENT_W, Inches(1.5), "Go-to-Market Slide Options", font_name=FONT_SANS, font_size=54, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN_X, Inches(4.2), CONTENT_W, Inches(0.5), "5 different approaches for reaching your first 40-60 beta users", font_name=FONT_BODY, font_size=20, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN_X, Inches(5.5), CONTENT_W, Inches(0.5), "Navigate through the options", font_name=FONT_BODY, font_size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ==============================================================================
# SLIDE 2: Option 1 - 4-Channel Beta Recruitment
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Go-to-Market / Option 1", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(1), "4-Channel Beta Recruitment Strategy", font_name=FONT_SANS, font_size=36, bold=True)

# Grid Layout
grid_y = Inches(2.2)
gap = Inches(0.3)
col_w = (CONTENT_W - (3 * gap)) / 4
box_h = Inches(3.5)

channels = [
    {
        "num": "01", "title": "MIT Campus", "target": "Target: 15-20 users", "color": ACCENT_GREEN,
        "items": ["Dorm workshops", "Student group demos", "Campus email blast", "Friend referrals"]
    },
    {
        "num": "02", "title": "Online Communities", "target": "Target: 15-20 users", "color": ACCENT_PURPLE,
        "items": ["r/productivity", "r/getdisciplined", "Indie Hackers", "Discord servers"]
    },
    {
        "num": "03", "title": "Content Marketing", "target": "Target: 10-15 users", "color": ACCENT_PURPLE,
        "items": ["2-3 TikToks/week", "Problem-solution posts", "Beta access CTA", "Behind-the-scenes"]
    },
    {
        "num": "04", "title": "Micro-Influencers", "target": "Target: 5-10 users", "color": TEXT_PRIMARY,
        "items": ["Reach out to 20", "Early access offer", "Feedback partnership", "Organic mentions"]
    }
]

for i, ch in enumerate(channels):
    x = MARGIN_X + i * (col_w + gap)
    fill = None
    if ch["color"] == ACCENT_GREEN:
        fill = RGBColor(20, 25, 0) # Dark green tint
    elif ch["color"] == ACCENT_PURPLE:
        fill = RGBColor(16, 14, 25) # Dark purple tint
    
    add_box(slide, x, grid_y, col_w, box_h, border=ch["color"], fill_color=fill, line_width=2)
    
    # Content
    add_textbox(slide, x + Inches(0.2), grid_y + Inches(0.2), col_w, Inches(0.5), ch["num"], font_name=FONT_SANS, font_size=32, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, x + Inches(0.2), grid_y + Inches(0.8), col_w - Inches(0.4), Inches(0.3), ch["title"].upper(), font_name=FONT_SANS, font_size=14, bold=True)
    add_textbox(slide, x + Inches(0.2), grid_y + Inches(1.2), col_w - Inches(0.4), Inches(0.3), ch["target"], font_size=11, bold=True, color=ACCENT_GREEN)
    add_bullets(slide, x + Inches(0.1), grid_y + Inches(1.5), col_w - Inches(0.2), Inches(1.5), ch["items"], font_size=11)

# Bottom Card
card_y = grid_y + box_h + Inches(0.4)
add_box(slide, MARGIN_X, card_y, CONTENT_W, Inches(0.8), border=ACCENT_GREEN, fill_color=RGBColor(20, 25, 0), line_width=2)
add_textbox(slide, MARGIN_X, card_y + Inches(0.2), CONTENT_W, Inches(0.4), "Parallel execution across 4 channels to reach 50+ beta users in 4-6 weeks", font_size=14, bold=True, align=PP_ALIGN.CENTER)


# ==============================================================================
# SLIDE 3: Option 2 - 3-Phase Launch
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Go-to-Market / Option 2", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(1), "3-Phase Beta Launch Strategy", font_name=FONT_SANS, font_size=36, bold=True)

# Timeline
tl_y = Inches(2.2)
phase_w = CONTENT_W / 3
phase_h = Inches(2.8)

phases = [
    {
        "num": "W1-2", "title": "Warm Network", "active": True,
        "items": ["MIT friends & classmates", "Personal social media", "Direct outreach", "Target: 15 users"]
    },
    {
        "num": "W3-4", "title": "Community Seeding", "active": False,
        "items": ["Reddit posts (5-8 subs)", "ProductHunt 'coming soon'", "Discord communities", "Target: 20-25 users"]
    },
    {
        "num": "W5-6", "title": "Content Amplification", "active": False,
        "items": ["12-15 content pieces", "Micro-influencer outreach", "User testimonials", "Target: 15-20 users"]
    }
]

for i, ph in enumerate(phases):
    x = MARGIN_X + i * phase_w
    border_color = ACCENT_GREEN if ph["active"] else TEXT_PRIMARY
    fill = RGBColor(20, 25, 0) if ph["active"] else None
    
    # Since they share borders in CSS, we can just overlap or butt them against each other.
    # Drawing simple boxes next to each other
    add_box(slide, x, tl_y, phase_w, phase_h, border=border_color, fill_color=fill, line_width=2)
    
    add_textbox(slide, x + Inches(0.2), tl_y + Inches(0.2), phase_w, Inches(0.6), ph["num"], font_name=FONT_SANS, font_size=32, bold=True, color=ACCENT_GREEN if ph["active"] else TEXT_MUTED)
    add_textbox(slide, x + Inches(0.2), tl_y + Inches(0.8), phase_w, Inches(0.3), ph["title"].upper(), font_name=FONT_SANS, font_size=14, bold=True)
    add_bullets(slide, x + Inches(0.1), tl_y + Inches(1.1), phase_w - Inches(0.2), Inches(1.5), ph["items"], font_size=12)

# Stats
stat_y = tl_y + phase_h + Inches(0.4)
stat_w = (CONTENT_W - Inches(0.4)) / 2
stat_h = Inches(1.5)

# Stat 1
add_box(slide, MARGIN_X, stat_y, stat_w, stat_h, border=ACCENT_GREEN, fill_color=RGBColor(20, 25, 0), line_width=2)
add_textbox(slide, MARGIN_X, stat_y + Inches(0.2), stat_w, Inches(0.6), "50-60", font_name=FONT_SANS, font_size=36, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN_X, stat_y + Inches(0.8), stat_w, Inches(0.3), "BETA USERS IN 6 WEEKS", font_size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# Stat 2
add_box(slide, MARGIN_X + stat_w + Inches(0.4), stat_y, stat_w, stat_h, border=TEXT_PRIMARY, line_width=2)
add_textbox(slide, MARGIN_X + stat_w + Inches(0.4), stat_y + Inches(0.2), stat_w, Inches(0.6), "$200-400", font_name=FONT_SANS, font_size=36, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN_X + stat_w + Inches(0.4), stat_y + Inches(0.8), stat_w, Inches(0.3), "ESTIMATED COST", font_size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# Note
note_y = stat_y + stat_h + Inches(0.2)
add_box(slide, MARGIN_X, note_y, CONTENT_W, Inches(0.6), border=TEXT_PRIMARY, line_width=1)
add_textbox(slide, MARGIN_X + Inches(0.2), note_y + Inches(0.15), CONTENT_W - Inches(0.4), Inches(0.3), "Sequential approach builds momentum and lets you refine messaging between phases.", font_size=12, color=TEXT_SECONDARY)


# ==============================================================================
# SLIDE 4: Option 3 - Problem-First
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Go-to-Market / Option 3", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(0.8), "Problem-First Community Approach", font_name=FONT_SANS, font_size=36, bold=True)
add_textbox(slide, MARGIN_X, Inches(1.4), CONTENT_W, Inches(0.5), "Lead with pain points, then offer solution access", font_size=16, color=TEXT_SECONDARY)

col_gap = Inches(0.6)
col_w = (CONTENT_W - col_gap) / 2
left_x = MARGIN_X
right_x = MARGIN_X + col_w + col_gap
content_y = Inches(2.2)

# Left Col
add_textbox(slide, left_x, content_y, col_w, Inches(0.4), "Where Your Users Are Struggling", font_name=FONT_SANS, font_size=16, bold=True, color=ACCENT_GREEN)

problems = [
    ("r/productivity (3.5M members)", "\"Can't stick to my routine\"",),
    ("r/getdisciplined (1.5M)", "\"Apps don't help me understand why\"",),
    ("r/selfimprovement (1.2M)", "\"Need science-backed approach\"",),
    ("TikTok comments", "\"How do I actually DO this?\"",)
]
prob_y = content_y + Inches(0.6)
for title, sub in problems:
    add_textbox(slide, left_x, prob_y, col_w, Inches(0.3), f"• {title}", font_size=14, bold=True)
    add_textbox(slide, left_x + Inches(0.3), prob_y + Inches(0.3), col_w, Inches(0.3), sub, font_size=12, color=TEXT_SECONDARY)
    prob_y += Inches(0.8)

# Right Col
add_textbox(slide, right_x, content_y, col_w, Inches(0.4), "Engagement Strategy", font_name=FONT_SANS, font_size=16, bold=True, color=ACCENT_GREEN)

cards = [
    ("Week 1-2: Listen & Engage", "Comment on frustration posts, offer insights, build credibility", ACCENT_PURPLE),
    ("Week 3-4: Share Solution", "Value-first posts about behavioral science + beta access offer", ACCENT_PURPLE),
    ("Week 5-6: Amplify", "User testimonials, insights from beta cohort", ACCENT_GREEN)
]
card_y = content_y + Inches(0.6)
for title, desc, color in cards:
    fill = RGBColor(20, 25, 0) if color == ACCENT_GREEN else None
    add_box(slide, right_x, card_y, col_w, Inches(1.2), border=color, fill_color=fill, line_width=2)
    add_textbox(slide, right_x + Inches(0.2), card_y + Inches(0.15), col_w - Inches(0.4), Inches(0.3), title, font_size=13, bold=True)
    add_textbox(slide, right_x + Inches(0.2), card_y + Inches(0.45), col_w - Inches(0.4), Inches(0.6), desc, font_size=12, color=TEXT_SECONDARY)
    card_y += Inches(1.4)

# Bottom "Why This Works"
why_y = Inches(6.5)
add_box(slide, MARGIN_X, why_y, CONTENT_W, Inches(0.8), border=TEXT_PRIMARY, line_width=1)
add_textbox(slide, MARGIN_X + Inches(0.2), why_y + Inches(0.1), CONTENT_W, Inches(0.3), "Why This Works", font_size=12, bold=True)
add_textbox(slide, MARGIN_X + Inches(0.2), why_y + Inches(0.35), CONTENT_W - Inches(0.4), Inches(0.4), "You're not selling an app - you're offering a solution to a problem they're already posting about daily.", font_size=12, color=TEXT_SECONDARY)


# ==============================================================================
# SLIDE 5: Option 4 - Influencer Model
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Go-to-Market / Option 4", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(0.8), "Influencer Micro-Partnership Model", font_name=FONT_SANS, font_size=36, bold=True)

left_x = MARGIN_X
right_x = MARGIN_X + col_w + col_gap
content_y = Inches(2.0)

# Left
# Stat Block
add_box(slide, left_x, content_y, col_w, Inches(1.2), border=ACCENT_GREEN, fill_color=RGBColor(20, 25, 0), line_width=2)
add_textbox(slide, left_x, content_y + Inches(0.2), col_w, Inches(0.5), "20", font_name=FONT_SANS, font_size=36, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_textbox(slide, left_x, content_y + Inches(0.7), col_w, Inches(0.3), "MICRO-INFLUENCERS TO REACH", font_size=10, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

add_textbox(slide, left_x, content_y + Inches(1.5), col_w, Inches(0.3), "Target Profile", font_name=FONT_SANS, font_size=14, bold=True)
add_bullets(slide, left_x, content_y + Inches(1.8), col_w, Inches(1.5), ["5K-50K followers", "Productivity/wellness niche", "High engagement rate (>3%)", "Creates routine content"])

plat_y = content_y + Inches(3.5)
add_box(slide, left_x, plat_y, col_w, Inches(0.8), border=TEXT_PRIMARY, line_width=1)
add_textbox(slide, left_x + Inches(0.2), plat_y + Inches(0.1), col_w, Inches(0.2), "Platforms to Target", font_size=11, bold=True)
add_textbox(slide, left_x + Inches(0.2), plat_y + Inches(0.35), col_w, Inches(0.3), "TikTok, Instagram, YouTube", font_size=11, color=TEXT_SECONDARY)

# Right - Partnership Offer
add_textbox(slide, right_x, content_y, col_w, Inches(0.4), "Partnership Offer", font_name=FONT_SANS, font_size=16, bold=True, color=ACCENT_GREEN)

offer_cards = [
    ("Early Access + Co-Creation", "Exclusive beta access + input on routine templates", ACCENT_PURPLE),
    ("Custom Routine Template", "Feature their routine as a template in the app", ACCENT_PURPLE),
    ("Lifetime Pro Access", "Free premium when you launch monetization", ACCENT_PURPLE),
    ("Expected Conversion", "20 outreach -> 5-7 partnerships -> 30-40 beta users", ACCENT_GREEN)
]
cy = content_y + Inches(0.5)
for title, desc, color in offer_cards:
    fill = RGBColor(20, 25, 0) if color == ACCENT_GREEN else None
    add_box(slide, right_x, cy, col_w, Inches(0.9), border=color, fill_color=fill, line_width=2)
    add_textbox(slide, right_x + Inches(0.15), cy + Inches(0.1), col_w - Inches(0.3), Inches(0.3), title, font_size=12, bold=True)
    add_textbox(slide, right_x + Inches(0.15), cy + Inches(0.4), col_w - Inches(0.3), Inches(0.4), desc, font_size=10, color=TEXT_SECONDARY)
    cy += Inches(1.05)

# Funnel Bottom
funnel_y = Inches(6.5)
funnel_h = Inches(0.8)
stages = [("Reach Out", "20"), ("Respond", "10-12"), ("Partner", "5-7"), ("Beta Users", "30-40")]
stage_w = CONTENT_W / 4
for i, (label, val) in enumerate(stages):
    sx = MARGIN_X + i * stage_w
    highlight = i == 3
    border = ACCENT_GREEN if highlight else TEXT_PRIMARY
    fill = RGBColor(20, 25, 0) if highlight else None
    
    add_box(slide, sx, funnel_y, stage_w, funnel_h, border=border, fill_color=fill, line_width=2)
    add_textbox(slide, sx + Inches(0.1), funnel_y + Inches(0.25), stage_w/2, Inches(0.3), label.upper(), font_size=10, bold=True)
    add_textbox(slide, sx + stage_w/2, funnel_y + Inches(0.15), stage_w/2 - Inches(0.1), Inches(0.5), val, font_name=FONT_SANS, font_size=20, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)


# ==============================================================================
# SLIDE 6: Option 5 - MIT Pipeline
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Go-to-Market / Option 5", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(0.8), "MIT-to-Market Pipeline", font_name=FONT_SANS, font_size=36, bold=True)
add_textbox(slide, MARGIN_X, Inches(1.4), CONTENT_W, Inches(0.5), "Validate on campus, then expand with proven model", font_size=16, color=TEXT_SECONDARY)

# 3 Columns
gap = Inches(0.4)
cw = (CONTENT_W - 2 * gap) / 3
cy = Inches(2.2)
ch = Inches(3.2)

cols = [
    {
        "title": "Phase 1: MIT Launch", "time": "Weeks 1-3", "color": ACCENT_GREEN,
        "items": ["Campus email blast", "Dorm presentations", "Student group demos", "Referral incentive"],
        "target": "Target: 25-30 users"
    },
    {
        "title": "Phase 2: Campus Expansion", "time": "Weeks 4-5", "color": ACCENT_PURPLE,
        "items": ["Harvard, BU, Northeastern", "Student ambassador program", "College subreddits", "Uni Discord servers"],
        "target": "Target: 15-20 users"
    },
    {
        "title": "Phase 3: Online Communities", "time": "Weeks 6+", "color": TEXT_PRIMARY,
        "items": ["Reddit productivity subs", "ProductHunt launch", "Indie Hackers", "Content marketing"],
        "target": "Target: 10-15 users"
    }
]

for i, col in enumerate(cols):
    cx = MARGIN_X + i * (cw + gap)
    fill = None
    if col["color"] == ACCENT_GREEN:
        fill = RGBColor(20, 25, 0)
    elif col["color"] == ACCENT_PURPLE:
        fill = RGBColor(16, 14, 25)
        
    add_box(slide, cx, cy, cw, ch, border=col["color"], fill_color=fill, line_width=2)
    
    # Header
    add_textbox(slide, cx + Inches(0.15), cy + Inches(0.2), cw - Inches(0.3), Inches(0.3), col["title"], font_size=13, bold=True)
    add_textbox(slide, cx + Inches(0.15), cy + Inches(0.45), cw - Inches(0.3), Inches(0.2), col["time"], font_size=11, color=TEXT_SECONDARY)
    
    # List
    add_bullets(slide, cx + Inches(0.1), cy + Inches(0.8), cw - Inches(0.2), Inches(1.5), col["items"], font_size=11)
    
    # Footer
    line_y = cy + ch - Inches(0.6)
    add_box(slide, cx + Inches(0.15), line_y, cw - Inches(0.3), Pt(1), fill_color=RGBColor(50, 50, 50))
    add_textbox(slide, cx + Inches(0.15), line_y + Inches(0.1), cw - Inches(0.3), Inches(0.3), col["target"], font_size=12, bold=True, color=col["color"])

# Bottom "Why Start with MIT"
why_y = cy + ch + Inches(0.4)
add_box(slide, MARGIN_X, why_y, CONTENT_W, Inches(0.8), border=ACCENT_GREEN, fill_color=RGBColor(20, 25, 0), line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), why_y + Inches(0.1), CONTENT_W, Inches(0.3), "Why Start with MIT", font_size=12, bold=True)
add_textbox(slide, MARGIN_X + Inches(0.2), why_y + Inches(0.35), CONTENT_W - Inches(0.4), Inches(0.4), "Credibility boost, easy access for interviews, homogeneous cohort for better insights.", font_size=12, color=TEXT_SECONDARY)


# ==============================================================================
# SLIDE 7: Summary
# ==============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide)

add_section_label(slide, "Summary", MARGIN_X, Inches(0.4))
add_textbox(slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(0.8), "Which Strategy Fits Best?", font_name=FONT_SANS, font_size=36, bold=True)

# Table-like structure
ty = Inches(2.0)
row_h = Inches(0.6)
cols_def = [
    ("Option", Inches(2.5)),
    ("Best For", Inches(3.5)),
    ("Time", Inches(1.5)),
    ("Budget", Inches(2.0)),
    ("Risk", Inches(1.5))
]
total_w = sum(c[1] for c in cols_def)
# Actually let's use exact table if possible, or just boxes. Boxes give us more style control matching previous slides.
# Headers
current_x = MARGIN_X
for name, w in cols_def:
    add_textbox(slide, current_x, ty, w, Inches(0.3), name, font_size=12, bold=True, color=TEXT_SECONDARY)
    current_x += w

# Rows
rows = [
    ["4-Channel Funnel", "Maximum reach & speed", "High", "$400-600", "Low"],
    ["3-Phase Launch", "Controlled growth & learning", "Medium", "$200-400", "Low"],
    ["Problem-First", "Organic engagement", "Medium", "$0-200", "Medium"],
    ["Micro-Influencers", "Fast user acquisition via trust", "Low-Medium", "$0 (equity)", "Med-High"],
    ["MIT-to-Market", "Credibility + Access (Recommended)", "Medium", "$300-500", "Low"]
]

ty += Inches(0.4)
for i, row in enumerate(rows):
    is_rec = i == 4
    fill = RGBColor(20, 25, 0) if is_rec else None
    
    # Row bg
    add_box(slide, MARGIN_X, ty, CONTENT_W, row_h, fill_color=fill)
    
    curr_x = MARGIN_X
    for j, cell_text in enumerate(row):
        w = cols_def[j][1]
        font_weight = True if j == 0 else False
        color = ACCENT_GREEN if is_rec and j == 0 else TEXT_PRIMARY
        add_textbox(slide, curr_x, ty + Inches(0.15), w, row_h, cell_text, font_size=12, bold=font_weight, color=color)
        curr_x += w
        
    # Border line
    add_box(slide, MARGIN_X, ty + row_h, CONTENT_W, Pt(1), fill_color=RGBColor(50, 50, 50))
    ty += row_h

# Recommendation Box
rec_y = ty + Inches(0.4)
add_box(slide, MARGIN_X, rec_y, CONTENT_W, Inches(1.2), border=ACCENT_GREEN, fill_color=RGBColor(20, 25, 0), line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), rec_y + Inches(0.2), CONTENT_W, Inches(0.3), "// RECOMMENDATION", font_name=FONT_MONO, font_size=12, bold=True, color=ACCENT_GREEN)
add_textbox(slide, MARGIN_X + Inches(0.2), rec_y + Inches(0.5), CONTENT_W - Inches(0.4), Inches(0.6), "Option 5 (MIT-to-Market) offers the best balance of credibility, access, and controlled expansion - while keeping costs within your budget.", font_size=16, bold=True)


# Save
output_path = "/workspaces/DMDHW/MIT_Sandbox/Pitch Deck/gtm-slide-options.pptx"
prs.save(output_path)
print(f"Saved PPTX to {output_path}")
