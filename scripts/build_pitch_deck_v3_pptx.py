from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0, 0, 0)
BG_SECONDARY = RGBColor(10, 10, 10)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(192, 192, 192)
TEXT_MUTED = RGBColor(102, 102, 102)
ACCENT_PURPLE = RGBColor(167, 139, 250)
ACCENT_GREEN = RGBColor(204, 255, 0)
ACCENT_RED = RGBColor(255, 0, 0)

# Fonts
FONT_SANS = "Arial Black"
FONT_BODY = "Arial"
FONT_MONO = "Consolas"

# Layout constants
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.55)
GAP_COL = Inches(0.6)
CONTENT_W = SLIDE_W - (MARGIN_X * 2)
COL_W = (CONTENT_W - GAP_COL) / 2


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, uppercase=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def add_section_label(slide, text, x, y):
    box = add_textbox(slide, x, y, Inches(4), Inches(0.3), f"// {text}", font_name=FONT_MONO, font_size=14, bold=True, color=ACCENT_PURPLE)
    return box


def add_h1(slide, text, x, y, w):
    return add_textbox(slide, x, y, w, Inches(0.9), text, font_name=FONT_SANS, font_size=42, bold=True, color=TEXT_PRIMARY)


def add_h3(slide, text, x, y, w, color=TEXT_PRIMARY):
    return add_textbox(slide, x, y, w, Inches(0.4), text, font_name=FONT_SANS, font_size=20, bold=True, color=color)


def add_body(slide, text, x, y, w, h, color=TEXT_SECONDARY, font_size=18, bold=False):
    return add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=font_size, bold=bold, color=color)


def add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=None, line_width=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.color.rgb = border
    shape.line.width = Pt(line_width)
    shape.fill.solid()
    if fill_color is None:
        shape.fill.fore_color.rgb = BG
    else:
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_logo(slide, x, y, size):
    # Draw a simplified version of the logo using shapes (white rounded square + black bars).
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    base.fill.solid()
    base.fill.fore_color.rgb = TEXT_PRIMARY
    base.line.color.rgb = TEXT_PRIMARY
    base.line.width = Pt(1)
    # Bars
    bar_h = size * 0.16
    bar_w1 = size * 0.7
    bar_w2 = size * 0.45
    bar_w3 = size * 0.18
    offset_x = x + size * 0.15
    base_y = y + size * 0.62
    mid_y = y + size * 0.42
    top_y = y + size * 0.22
    bars = [
        (offset_x, base_y, bar_w1, bar_h),
        (x + size * 0.28, mid_y, bar_w2, bar_h),
        (x + size * 0.41, top_y, bar_w3, bar_h),
    ]
    for bx, by, bw, bh in bars:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bw, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = BG
        bar.line.fill.background()


def add_bullets(slide, x, y, w, h, items, font_size=18, color=TEXT_PRIMARY, bullet_color=ACCENT_GREEN):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.space_before = Pt(6)
    return box


def add_stat_block(slide, x, y, w, h, value, label, color=ACCENT_GREEN, bg_alpha=0.05):
    box = add_box(slide, x, y, w, h, border=color, fill_color=RGBColor(10, 10, 10), line_width=2.5)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = value
    run1.font.name = FONT_SANS
    run1.font.size = Pt(40)
    run1.font.bold = True
    run1.font.color.rgb = color
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER
    return box


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
logo_size = Inches(1.1)
add_logo(slide, (SLIDE_W - logo_size) / 2, Inches(1.0), logo_size)
add_textbox(slide, Inches(0.5), Inches(2.1), SLIDE_W - Inches(1.0), Inches(0.8), "Axiom Forge", font_name=FONT_SANS, font_size=56, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(3.0), SLIDE_W - Inches(4.0), Inches(0.35), "PRINCIPLES TURNED INTO PRACTICE", font_name=FONT_BODY, font_size=20, bold=False, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
# Highlight line
highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(3.45), SLIDE_W - Inches(4.4), Inches(0.45))
highlight.fill.solid()
highlight.fill.fore_color.rgb = ACCENT_GREEN
highlight.line.fill.background()
add_textbox(slide, Inches(2.2), Inches(3.45), SLIDE_W - Inches(4.4), Inches(0.45), "ROUTINES THAT TURN INTENT INTO IDENTITY", font_name=FONT_BODY, font_size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
add_section_label(slide, "MIT Sandbox Pitch Deck", Inches(4.3), Inches(4.25))
# Footer meta boxes
meta_y = Inches(6.3)
meta_w = Inches(3.3)
meta_h = Inches(0.45)
meta_xs = [Inches(2.1), Inches(5.1), Inches(8.1)]
meta_text = ["Pitch Deck for Working Prototype", "January 2026", "Javier Serrano"]
for i, mx in enumerate(meta_xs):
    box = add_box(slide, mx, meta_y, meta_w, meta_h, border=TEXT_MUTED, fill_color=BG, line_width=1)
    tf = box.text_frame
    tf.text = meta_text[i]
    tf.paragraphs[0].font.name = FONT_BODY
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = TEXT_SECONDARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 2: Context
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Context", MARGIN_X, MARGIN_TOP)
add_h1(slide, "Routines are the new obsession", MARGIN_X, Inches(0.95), CONTENT_W)
# Left column
left_x = MARGIN_X
left_y = Inches(2.0)
add_stat_block(slide, left_x, left_y, COL_W, Inches(1.2), "+50B", "morning routine videos on TikTok")
# Hashtag pills
pill_y = left_y + Inches(1.35)
for idx, (name, views) in enumerate([("#ThatGirl", "17.4B views"), ("#nightroutine", "11.4B views"), ("#dailyroutine", "8B+ views")]):
    pill = add_box(slide, left_x + (idx % 2) * Inches(2.9), pill_y + (idx // 2) * Inches(0.55), Inches(2.7), Inches(0.4), border=ACCENT_GREEN, fill_color=BG, line_width=1)
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.name = FONT_BODY
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_GREEN
    p2 = tf.add_paragraph()
    p2.text = views
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SECONDARY
# Chart container
chart_y = pill_y + Inches(1.3)
chart = add_box(slide, left_x, chart_y, COL_W, Inches(1.6), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_h3(slide, "Engagement vs. Other Content", left_x + Inches(0.2), chart_y + Inches(0.1), COL_W - Inches(0.4), color=ACCENT_GREEN)
# Bars
bar1_y = chart_y + Inches(0.55)
bar2_y = chart_y + Inches(1.05)
add_textbox(slide, left_x + Inches(0.2), bar1_y, Inches(1.2), Inches(0.3), "Routine", font_name=FONT_BODY, font_size=12, bold=True, color=TEXT_PRIMARY)
bar1 = add_box(slide, left_x + Inches(1.6), bar1_y, Inches(2.8), Inches(0.25), border=TEXT_PRIMARY, fill_color=BG_SECONDARY, line_width=1)
fill1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x + Inches(1.6), bar1_y, Inches(2.8), Inches(0.25))
fill1.fill.solid()
fill1.fill.fore_color.rgb = ACCENT_GREEN
fill1.line.fill.background()
add_textbox(slide, left_x + Inches(4.6), bar1_y, Inches(0.6), Inches(0.3), "6-10x", font_name=FONT_BODY, font_size=12, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)
add_textbox(slide, left_x + Inches(0.2), bar2_y, Inches(1.2), Inches(0.3), "General Wellness", font_name=FONT_BODY, font_size=11, bold=True, color=TEXT_PRIMARY)
bar2 = add_box(slide, left_x + Inches(1.6), bar2_y, Inches(2.8), Inches(0.25), border=TEXT_PRIMARY, fill_color=BG_SECONDARY, line_width=1)
fill2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x + Inches(1.6), bar2_y, Inches(0.5), Inches(0.25))
fill2.fill.solid()
fill2.fill.fore_color.rgb = TEXT_MUTED
fill2.line.fill.background()
add_textbox(slide, left_x + Inches(4.6), bar2_y, Inches(0.6), Inches(0.3), "1x", font_name=FONT_BODY, font_size=12, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.RIGHT)
# Right column influencers
right_x = MARGIN_X + COL_W + GAP_COL
add_h3(slide, "Influencers Proving Demand", right_x, Inches(2.0), COL_W, color=ACCENT_GREEN)
items = [
    "Andrew Huberman\n7.3M subs - Neuroscience-backed routines dominate",
    "Ashton Hall\n3M+ followers - Routine videos avg 1M+ views\n+750 Million views on X",
    "Ali Abdaal\n6.5M subs - Productivity systems & habit stacking",
    "James Clear\n25M+ Atomic Habits sold - Created demand for systems",
]
inf_y = Inches(2.55)
for i, item in enumerate(items):
    box = slide.shapes.add_textbox(right_x, inf_y + Inches(0.85) * i, COL_W, Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    title, *rest = item.split("\n")
    run.text = title
    run.font.name = FONT_BODY
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT_PRIMARY
    for line in rest:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = ACCENT_GREEN if "views" in line else TEXT_SECONDARY
# Highlight card
highlight = add_box(slide, MARGIN_X, Inches(6.0), CONTENT_W, Inches(0.65), border=ACCENT_GREEN, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X, Inches(6.05), CONTENT_W, Inches(0.5), "Massive audience wants to DO routines, not just WATCH them.", font_name=FONT_BODY, font_size=18, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
# Sources
add_textbox(slide, MARGIN_X, Inches(6.75), CONTENT_W, Inches(0.3), "Sources: Broadcasting your breakfast: why TikTokers obsess over morning routines - The Guardian; Why We're Obsessed with Other People's Morning Routines - Time", font_name=FONT_BODY, font_size=12, color=TEXT_PRIMARY)

# Slide 3: Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Problem", MARGIN_X, MARGIN_TOP)
add_h1(slide, "But few people can execute them", MARGIN_X, Inches(0.95), CONTENT_W)
left_x = MARGIN_X
left_y = Inches(2.0)
metric_h = Inches(1.1)
metrics = [("66%", "of Gen Z use digital wellness tools"), ("~4%", "Day 30 retention rate in health & fitness apps (lowest category)"), ("~90%", "of people quit or fail when trying to form new habits")]
for i, (val, label) in enumerate(metrics):
    box = add_box(slide, left_x, left_y + i * Inches(1.25), COL_W, metric_h, border=ACCENT_GREEN, fill_color=BG, line_width=2)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = val
    p1.font.name = FONT_SANS
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GREEN
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER
# Right column - reasons
right_x = MARGIN_X + COL_W + GAP_COL
add_h3(slide, "Why People Fail", right_x, Inches(2.0), COL_W, color=ACCENT_RED)
reasons = [
    ("No Bridge from Content to Action", "Videos inspire but offer no execution tool"),
    ("Blank Slate Apps", "\"What do you want to track?\" - Users don't know"),
    ("No Science, No System", "Self-designed habits lack proven structure"),
    ("Streak Anxiety", "One miss = total abandonment"),
]
for i, (title, desc) in enumerate(reasons):
    box = slide.shapes.add_textbox(right_x, Inches(2.5) + Inches(0.75) * i, COL_W, Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = FONT_BODY
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_PRIMARY
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = FONT_BODY
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
# Opportunity full width
opp = add_box(slide, MARGIN_X, Inches(6.0), CONTENT_W, Inches(0.7), border=ACCENT_GREEN, fill_color=BG, line_width=2)
add_section_label(slide, "The Opportunity", MARGIN_X + Inches(0.2), Inches(6.05))
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(6.32), CONTENT_W - Inches(0.4), Inches(0.3), 'Users are truly struggling bridging from "routine inspiration" to "routine execution".', font_name=FONT_BODY, font_size=16, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, MARGIN_X, Inches(6.75), CONTENT_W, Inches(0.3), "Sources: The $2 trillion global wellness market gets a millennial and Gen Z glow-up - McKinsey & Company; Retention Rates for Mobile Apps by Industry - Plotline; Only 10% of people who try to create a habit achieve their goal - Summa Magazine", font_name=FONT_BODY, font_size=12, color=TEXT_PRIMARY)

# Slide 4: Solution
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Solution", MARGIN_X, MARGIN_TOP)
add_h1(slide, "Axiom Forge: A Habit System Built on Evidence, Reinforced by Mindset and Insight", MARGIN_X, Inches(0.95), CONTENT_W)
card_w = (CONTENT_W - Inches(0.6) * 2) / 3
card_h = Inches(4.4)
card_y = Inches(2.0)
card_xs = [MARGIN_X, MARGIN_X + card_w + Inches(0.6), MARGIN_X + (card_w + Inches(0.6)) * 2]
card_titles = ["Smart Routines", "Behavioral Insights", "Mindset Pills"]
card_subs = [
    "Science-backed templates remove the blank slate.",
    "Patterns that explain success or failure, generated from real data.",
    "Daily wisdom that reinforces identity.",
]
card_lists = [
    ["Morning + evening routines ready from day one", "Evidence-based structure for sleep, fitness, focus", "Designed for real schedules, not perfection"],
    ["Correlation, weekday patterns, and trend signals", "Sequence analysis for \"what works together\"", "Anomalies detection, strength scores and alerts"],
    ["Daily quote prompts from thought leaders and top performers", "Reflection journaling", "Rating & archive history"],
]
for i in range(3):
    border_color = ACCENT_PURPLE if i == 2 else TEXT_PRIMARY
    card = add_box(slide, card_xs[i], card_y, card_w, card_h, border=border_color, fill_color=BG, line_width=2)
    add_h3(slide, card_titles[i], card_xs[i] + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), color=TEXT_PRIMARY)
    add_body(slide, card_subs[i], card_xs[i] + Inches(0.2), card_y + Inches(0.7), card_w - Inches(0.4), Inches(0.5), color=TEXT_SECONDARY, font_size=12)
    list_y = card_y + Inches(1.2)
    add_bullets(slide, card_xs[i] + Inches(0.2), list_y, card_w - Inches(0.4), card_h - Inches(1.4), card_lists[i], font_size=12, color=TEXT_PRIMARY)

# Slide 5: Market
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Market", MARGIN_X, MARGIN_TOP)
add_h1(slide, "Axiom Forge taps into the rapidly growing Habit Tracking App Market, which is projected to triple within the next 10 years.", MARGIN_X, Inches(0.95), CONTENT_W)
chart_x = MARGIN_X
chart_y = Inches(2.1)
chart_w = CONTENT_W * 0.6
chart_h = Inches(4.2)
add_box(slide, chart_x, chart_y, chart_w, chart_h, border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, chart_x + Inches(0.2), chart_y + Inches(0.1), chart_w - Inches(0.4), Inches(0.3), "Global Habit Tracking App Market (USD Bn)", font_name=FONT_BODY, font_size=12, bold=True, color=ACCENT_GREEN)
# Bars
bars = [
    ("2025", "13.1", 0.25),
    ("2026", "14.9", 0.32),
    ("2027", "17.1", 0.38),
    ("2028", "", 0.44),
    ("2029", "", 0.51),
    ("2030", "", 0.57),
    ("2031", "", 0.65),
    ("2032", "", 0.74),
    ("2033", "", 0.84),
    ("2034", "", 0.94),
    ("2035", "50.2", 1.08),
]
bar_area_y = chart_y + Inches(0.6)
bar_area_h = chart_h - Inches(1.2)
bar_w = (chart_w - Inches(1.0)) / len(bars)
for i, (year, val, scale) in enumerate(bars):
    bx = chart_x + Inches(0.5) + bar_w * i
    bh = bar_area_h * scale
    by = bar_area_y + (bar_area_h - bh)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bar_w * 0.5, bh)
    bar.fill.solid()
    if year in ("2025", "2026", "2027"):
        bar.fill.fore_color.rgb = ACCENT_GREEN
        bar.line.fill.background()
    else:
        bar.fill.fore_color.rgb = RGBColor(100, 120, 0)
        bar.line.color.rgb = ACCENT_GREEN
        bar.line.width = Pt(1)
    if val:
        add_textbox(slide, bx - Inches(0.05), by - Inches(0.25), bar_w, Inches(0.2), val, font_name=FONT_BODY, font_size=10, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, bx - Inches(0.05), bar_area_y + bar_area_h + Inches(0.1), bar_w, Inches(0.2), year, font_name=FONT_BODY, font_size=9, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
# Diagonal CAGR arrows
line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_x + Inches(1.0), chart_y + Inches(0.55), chart_x + Inches(2.6), chart_y + Inches(0.35))
line1.line.color.rgb = ACCENT_PURPLE
line1.line.width = Pt(2)
line1.line.end_arrowhead = True
add_textbox(slide, chart_x + Inches(0.9), chart_y + Inches(0.2), Inches(1.5), Inches(0.3), "CAGR 14.3%", font_name=FONT_BODY, font_size=10, bold=True, color=ACCENT_PURPLE)
line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_x + Inches(2.6), chart_y + Inches(0.55), chart_x + Inches(5.6), chart_y + Inches(0.2))
line2.line.color.rgb = ACCENT_PURPLE
line2.line.width = Pt(2)
line2.line.end_arrowhead = True
add_textbox(slide, chart_x + Inches(2.7), chart_y + Inches(0.05), Inches(2.5), Inches(0.3), "CAGR 14.4%", font_name=FONT_BODY, font_size=10, bold=True, color=ACCENT_PURPLE)
# Drivers box
drivers_x = chart_x + chart_w + Inches(0.5)
drivers_w = CONTENT_W - chart_w - Inches(0.5)
drivers = add_box(slide, drivers_x, chart_y, drivers_w, chart_h, border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_h3(slide, "Market Growth Drivers", drivers_x + Inches(0.2), chart_y + Inches(0.2), drivers_w - Inches(0.4), color=ACCENT_GREEN)
items = [
    ("Rising consumer emphasis on self-improvement, mental well-being, and productivity enhancement", "68% of smartphone users engage with at least one productivity or wellness app"),
    ("Higher willingness to pay", "Premium subscriptions account for 36% of overall habit tracking app downloads"),
    ("Complementary to surge of wearables through integration", "Wearable integration improved user retention by 39%"),
]
cur_y = chart_y + Inches(0.8)
for title, stat in items:
    add_textbox(slide, drivers_x + Inches(0.3), cur_y, drivers_w - Inches(0.6), Inches(0.5), title, font_name=FONT_BODY, font_size=14, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, drivers_x + Inches(0.3), cur_y + Inches(0.45), drivers_w - Inches(0.6), Inches(0.4), stat, font_name=FONT_BODY, font_size=12, bold=True, color=ACCENT_GREEN)
    cur_y += Inches(1.2)
add_textbox(slide, MARGIN_X, Inches(6.75), CONTENT_W, Inches(0.3), "Source: Global Habit Tracking App Market Size - Global Growth Insights", font_name=FONT_BODY, font_size=12, color=TEXT_PRIMARY)

# Slide 6: Competition table
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Competition", MARGIN_X, MARGIN_TOP)
add_h1(slide, "Axiom Forge Is the Only App Combining Science + Mindset + Insights", MARGIN_X, Inches(0.95), CONTENT_W)
rows = 6
cols = 5
table = slide.shapes.add_table(rows, cols, MARGIN_X, Inches(2.1), CONTENT_W, Inches(3.0)).table
headers = ["App", "Science-Backed Templates", "Behavioral Insights", "Flexible Streaks", "Cross-Platform"]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.name = FONT_BODY
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = BG
    cell.fill.solid()
    cell.fill.fore_color.rgb = TEXT_PRIMARY
entries = [
    ["Streaks", "-", "-", "-", "iOS Only"],
    ["Habitica", "Generic", "Basic Stats", "Rest Inn", "Yes"],
    ["Loop", "-", "Good", "Yes", "Android"],
    ["Habitify", "-", "Mood Only", "-", "Yes"],
    ["Axiom Forge", "Evidence-Based", "Full Analysis", "Flexible", "PWA"],
]
for r, row in enumerate(entries, start=1):
    for c, value in enumerate(row):
        cell = table.cell(r, c)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_PRIMARY
        if r == len(entries):
            p.font.color.rgb = ACCENT_PURPLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG if r % 2 == 1 else BG_SECONDARY
# Positioning card
card = add_box(slide, MARGIN_X, Inches(5.4), CONTENT_W, Inches(0.9), border=ACCENT_GREEN, fill_color=BG, line_width=2)
add_section_label(slide, "Positioning", MARGIN_X + Inches(0.2), Inches(5.45))
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(5.75), CONTENT_W - Inches(0.4), Inches(0.4), "The only habit tracker that pairs science-backed routines with automated behavioral insight.", font_name=FONT_BODY, font_size=16, bold=True, color=TEXT_PRIMARY)

# Slide 7: Prototype
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Prototype", MARGIN_X, MARGIN_TOP)
add_h1(slide, "A Working Product Validated by Early Users", MARGIN_X, Inches(0.95), CONTENT_W)
left_x = MARGIN_X
right_x = MARGIN_X + COL_W + GAP_COL
add_h3(slide, "Current State", left_x, Inches(2.0), COL_W)
add_bullets(slide, left_x, Inches(2.5), COL_W, Inches(3.0), [
    "PWA live with Firebase auth + sync",
    "Morning & evening routines with daily check-ins",
    "Dashboard with streaks, heatmap, and completion rates",
    "Smart Insights engine (correlation, trends, anomalies)",
    "Mindset journal with quotes + reflections",
], font_size=14, color=TEXT_PRIMARY)
add_h3(slide, "Validation", right_x, Inches(2.0), COL_W)
card1 = add_box(slide, right_x, Inches(2.5), COL_W, Inches(1.2), border=ACCENT_GREEN, fill_color=BG, line_width=2)
add_textbox(slide, right_x + Inches(0.2), Inches(2.65), COL_W - Inches(0.4), Inches(0.4), "Validated with several early users", font_name=FONT_BODY, font_size=16, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, right_x + Inches(0.2), Inches(3.05), COL_W - Inches(0.4), Inches(0.4), "Qualitative feedback confirms clarity, motivation, and desire for deeper insights.", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)
card2 = add_box(slide, right_x, Inches(4.0), COL_W, Inches(1.1), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, right_x + Inches(0.2), Inches(4.15), COL_W - Inches(0.4), Inches(0.4), "Next Validation Goal", font_name=FONT_BODY, font_size=14, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, right_x + Inches(0.2), Inches(4.55), COL_W - Inches(0.4), Inches(0.4), "Run 40-60 user interviews + beta cohort to confirm retention and pricing.", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)

# Slide 8: Business Model
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Business Model", MARGIN_X, MARGIN_TOP)
add_h1(slide, "Multiple Revenue Streams from Freemium to Enterprise", MARGIN_X, Inches(0.95), CONTENT_W)
left_x = MARGIN_X
right_x = MARGIN_X + COL_W + GAP_COL
cards_left = [
    ("FREEMIUM CORE", "Tracking, streaks, routines, basic analytics.", TEXT_PRIMARY, None),
    ("INSIGHTS PRO", "Correlations, trends, sequences, advanced dashboards.", ACCENT_PURPLE, None),
    ("COACH LAYER", "AI summaries + habit recommendations.", ACCENT_GREEN, None),
]
for i, (title, desc, border, _) in enumerate(cards_left):
    card = add_box(slide, left_x, Inches(2.1) + Inches(1.2) * i, COL_W, Inches(1.0), border=border, fill_color=BG, line_width=2)
    add_textbox(slide, left_x + Inches(0.2), Inches(2.2) + Inches(1.2) * i, COL_W - Inches(0.4), Inches(0.4), title, font_name=FONT_BODY, font_size=16, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, left_x + Inches(0.2), Inches(2.6) + Inches(1.2) * i, COL_W - Inches(0.4), Inches(0.4), desc, font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)
cards_right = [
    ("CONTENT PACKS", "Premium routine templates and expert programs."),
    ("TEAMS / CAMPUS", "Group dashboards for student cohorts or clubs."),
    ("LIFETIME LICENSE", "One-time fee for pro features."),
]
for i, (title, desc) in enumerate(cards_right):
    card = add_box(slide, right_x, Inches(2.1) + Inches(1.2) * i, COL_W, Inches(1.0), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
    add_textbox(slide, right_x + Inches(0.2), Inches(2.2) + Inches(1.2) * i, COL_W - Inches(0.4), Inches(0.4), title, font_name=FONT_BODY, font_size=16, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, right_x + Inches(0.2), Inches(2.6) + Inches(1.2) * i, COL_W - Inches(0.4), Inches(0.4), desc, font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)

# Slide 9: Budget
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Funding Request", MARGIN_X, MARGIN_TOP)
add_h1(slide, "$4,500 to Validate, Polish, and Launch to Market", MARGIN_X, Inches(0.95), CONTENT_W)
wrap = add_box(slide, MARGIN_X, Inches(2.0), CONTENT_W, Inches(3.6), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(2.15), Inches(2.0), Inches(0.6), "$4,500", font_name=FONT_SANS, font_size=40, bold=True, color=ACCENT_GREEN)
add_textbox(slide, MARGIN_X + Inches(2.4), Inches(2.3), Inches(7.5), Inches(0.4), "Aligned with MIT Sandbox reimbursable categories", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)
# Budget bars
items = [
    ("User Research", 1200, ACCENT_GREEN),
    ("Contractors", 1000, ACCENT_PURPLE),
    ("Software & Tools", 900, ACCENT_PURPLE),
    ("Marketing & Launch", 800, ACCENT_PURPLE),
    ("Legal & Compliance", 600, ACCENT_PURPLE),
]
bar_y = Inches(3.0)
for i, (label, amt, color) in enumerate(items):
    y = bar_y + Inches(0.5) * i
    add_textbox(slide, MARGIN_X + Inches(0.2), y, Inches(1.8), Inches(0.3), label.upper(), font_name=FONT_BODY, font_size=11, bold=True, color=TEXT_PRIMARY)
    bar = add_box(slide, MARGIN_X + Inches(2.2), y, Inches(6.5), Inches(0.25), border=TEXT_PRIMARY, fill_color=BG_SECONDARY, line_width=1)
    fill_w = Inches(6.5 * (amt / 4500))
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X + Inches(2.2), y, fill_w, Inches(0.25))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    add_textbox(slide, MARGIN_X + Inches(9.0), y, Inches(1.0), Inches(0.3), f"${amt}", font_name=FONT_BODY, font_size=11, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.RIGHT)
# Note card
note = add_box(slide, MARGIN_X, Inches(5.8), CONTENT_W, Inches(0.7), border=TEXT_PRIMARY, fill_color=BG, line_width=1)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(5.95), CONTENT_W - Inches(0.4), Inches(0.3), "Contractor spend is within the 25% cap. Research incentives and software licenses follow Sandbox reimbursement guidelines.", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)

# Slide 10: Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Roadmap", MARGIN_X, MARGIN_TOP)
add_h1(slide, "90-Day Sprint from Validation to App Store Launch", MARGIN_X, Inches(0.95), CONTENT_W)
phase_w = (CONTENT_W) / 3
phase_y = Inches(2.1)
phase_h = Inches(3.0)
phases = [
    ("01", "Validate", ["40-60 user interviews", "Beta cohort with retention tracking", "Pricing tests for Insights Pro"], True),
    ("02", "Polish", ["Refine insights UI", "Performance + onboarding improvements", "Brand + marketing assets"], False),
    ("03", "Launch", ["App Store packaging", "Public release + press outreach", "Post-launch analytics + iteration"], False),
]
for i, (num, title, bullets, active) in enumerate(phases):
    px = MARGIN_X + phase_w * i
    border = ACCENT_GREEN if active else TEXT_PRIMARY
    card = add_box(slide, px, phase_y, phase_w, phase_h, border=border, fill_color=BG, line_width=2)
    add_textbox(slide, px + Inches(0.2), phase_y + Inches(0.1), Inches(1.0), Inches(0.5), num, font_name=FONT_SANS, font_size=36, bold=True, color=ACCENT_GREEN if active else TEXT_MUTED)
    add_textbox(slide, px + Inches(0.2), phase_y + Inches(0.7), phase_w - Inches(0.4), Inches(0.3), title, font_name=FONT_BODY, font_size=14, bold=True, color=TEXT_PRIMARY)
    add_bullets(slide, px + Inches(0.2), phase_y + Inches(1.1), phase_w - Inches(0.4), Inches(1.8), bullets, font_size=12, color=TEXT_SECONDARY)
# Launch target stat
stat = add_box(slide, MARGIN_X, Inches(5.5), Inches(3.5), Inches(0.9), border=ACCENT_GREEN, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(5.6), Inches(3.1), Inches(0.4), "Q1 2026", font_name=FONT_SANS, font_size=24, bold=True, color=ACCENT_GREEN)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(6.0), Inches(3.1), Inches(0.3), "App Store Launch Target", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)

# Slide 11: Team
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_section_label(slide, "Team", MARGIN_X, MARGIN_TOP)
add_h1(slide, "The Team Behind Axiom Forge", MARGIN_X, Inches(0.95), CONTENT_W)
card = add_box(slide, MARGIN_X, Inches(2.2), CONTENT_W, Inches(3.5), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X, Inches(3.7), CONTENT_W, Inches(0.5), "Team information coming soon...", font_name=FONT_BODY, font_size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

prs.save("/workspaces/DMDHW/MIT_Sandbox/Pitch Deck/axiom-forge-pitch-deck-v3.pptx")
print("Saved v3 PPTX")

# Slide Options v3: Extracted slides 2, 3, 6, 8 from slide-options-v3.html
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.55)
GAP_COL = Inches(0.6)
CONTENT_W = SLIDE_W - (MARGIN_X * 2)
COL_W = (CONTENT_W - GAP_COL) / 2


def add_cell(slide, x, y, w, h, text, fill_color=BG, text_color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.CENTER, font_size=12):
    cell = add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=fill_color, line_width=1)
    tf = cell.text_frame
    tf.clear()
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return cell


def add_option_label(slide, text, x, y):
    label_w = Inches(7.2)
    label_h = Inches(0.3)
    label = add_box(slide, x, y, label_w, label_h, border=ACCENT_PURPLE, fill_color=ACCENT_PURPLE, line_width=1)
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BG
    p.alignment = PP_ALIGN.LEFT
    return label


def add_phone_mockup(slide, x, y, w, h, border_color, play_color):
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    phone.fill.solid()
    phone.fill.fore_color.rgb = BG
    phone.line.color.rgb = border_color
    phone.line.width = Pt(2)
    screen = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + w * 0.08, y + h * 0.08, w * 0.84, h * 0.84)
    screen.fill.solid()
    screen.fill.fore_color.rgb = BG_SECONDARY
    screen.line.fill.background()
    notch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + w * 0.35, y + h * 0.1, w * 0.3, h * 0.05)
    notch.fill.solid()
    notch.fill.fore_color.rgb = BG
    notch.line.fill.background()
    play_size = w * 0.28
    play_x = x + (w - play_size) / 2
    play_y = y + (h - play_size) / 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, play_x, play_y, play_size, play_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG
    circle.line.color.rgb = play_color
    circle.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, play_x + play_size * 0.38, play_y + play_size * 0.28, play_size * 0.28, play_size * 0.32)
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = play_color
    tri.line.fill.background()
    return screen


# Slide 2: Option B - Tiered Layout
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_option_label(slide, "Slide 8 - Option B: Tiered Layout", MARGIN_X, Inches(0.35))
add_section_label(slide, "Business Model", MARGIN_X, Inches(0.8))
add_h1(slide, "Multiple Revenue Streams from Freemium to Enterprise", MARGIN_X, Inches(1.3), CONTENT_W)

base_y = Inches(2.2)
base_h = Inches(1.35)
base = add_box(slide, MARGIN_X, base_y, CONTENT_W, base_h, border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), base_y + Inches(0.1), CONTENT_W - Inches(0.4), Inches(0.3), "Freemium Core", font_name=FONT_SANS, font_size=18, bold=True, color=ACCENT_GREEN)
add_body(slide, "Everything you need to build habits: Tracking, streaks, morning & evening routines, basic analytics, correlations, trends, sequences, and advanced dashboards.", MARGIN_X + Inches(0.2), base_y + Inches(0.45), CONTENT_W - Inches(0.4), Inches(0.5), color=TEXT_SECONDARY, font_size=12)
add_textbox(slide, MARGIN_X + Inches(0.2), base_y + Inches(0.95), CONTENT_W - Inches(0.4), Inches(0.3), "FREE FOREVER", font_name=FONT_BODY, font_size=14, bold=True, color=ACCENT_GREEN)

tier_y = base_y + base_h + Inches(0.35)
tier_gap = Inches(0.25)
tier_w = (CONTENT_W - tier_gap * 3) / 4
tier_h = Inches(2.2)
tiers = [
    ("Coach Layer", "AI-powered summaries and personalized habit recommendations", "$4.99/mo", True),
    ("Content Packs", "Premium templates from top routine influencers", "$2.99/pack", False),
    ("Teams / Campus", "Group dashboards for cohorts and clubs", "$9.99/mo", False),
    ("Lifetime License", "Permanent access to all pro features", "$49.99", False),
]
for i, (title, desc, price, featured) in enumerate(tiers):
    x = MARGIN_X + i * (tier_w + tier_gap)
    border = ACCENT_PURPLE if featured else TEXT_PRIMARY
    fill = BG_SECONDARY if featured else BG
    card = add_box(slide, x, tier_y, tier_w, tier_h, border=border, fill_color=fill, line_width=2)
    add_textbox(slide, x + Inches(0.15), tier_y + Inches(0.12), tier_w - Inches(0.3), Inches(0.3), title, font_name=FONT_BODY, font_size=13, bold=True, color=ACCENT_PURPLE if featured else TEXT_PRIMARY)
    add_body(slide, desc, x + Inches(0.15), tier_y + Inches(0.45), tier_w - Inches(0.3), Inches(1.0), color=TEXT_SECONDARY, font_size=10)
    add_textbox(slide, x + Inches(0.15), tier_y + Inches(1.65), tier_w - Inches(0.3), Inches(0.3), price, font_name=FONT_BODY, font_size=14, bold=True, color=TEXT_PRIMARY)

# Slide 3: Option C - Pricing Grid
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_option_label(slide, "Slide 8 - Option C: Pricing Grid", MARGIN_X, Inches(0.35))
add_section_label(slide, "Business Model", MARGIN_X, Inches(0.8))
add_h1(slide, "Multiple Revenue Streams from Freemium to Enterprise", MARGIN_X, Inches(1.3), CONTENT_W)

table_y = Inches(2.1)
row_h = Inches(0.35)
feature_w = Inches(3.6)
col_w = (CONTENT_W - feature_w) / 4

headers = ["Features", "Freemium", "Coach Layer", "Content Packs", "Teams"]
for i, header in enumerate(headers):
    x = MARGIN_X + (feature_w if i > 0 else 0) + (col_w * (i - 1) if i > 1 else 0)
    w = feature_w if i == 0 else col_w
    fill = BG
    color = TEXT_PRIMARY
    if header == "Freemium":
        fill = ACCENT_GREEN
        color = BG
    elif header == "Coach Layer":
        fill = ACCENT_PURPLE
    add_cell(slide, x, table_y, w, row_h, header, fill_color=fill, text_color=color, bold=True, align=PP_ALIGN.CENTER, font_size=12)

rows = [
    ("Habit Tracking & Streaks", ["✓", "✓", "✓", "✓"]),
    ("Morning & Evening Routines", ["✓", "✓", "✓", "✓"]),
    ("Correlations & Trends", ["✓", "✓", "✓", "✓"]),
    ("Advanced Dashboards", ["✓", "✓", "✓", "✓"]),
    ("AI Habit Recommendations", ["-", "✓", "-", "-"]),
    ("Premium Routine Templates", ["-", "-", "✓", "-"]),
    ("Influencer Expert Programs", ["-", "-", "✓", "-"]),
    ("Group Dashboards", ["-", "-", "-", "✓"]),
]
for r_idx, (feature, values) in enumerate(rows, start=1):
    y = table_y + row_h * r_idx
    add_cell(slide, MARGIN_X, y, feature_w, row_h, feature, fill_color=BG, text_color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, font_size=11)
    for c_idx, value in enumerate(values):
        x = MARGIN_X + feature_w + col_w * c_idx
        color = ACCENT_GREEN if value == "✓" else TEXT_MUTED
        add_cell(slide, x, y, col_w, row_h, value, fill_color=BG, text_color=color, bold=True, align=PP_ALIGN.CENTER, font_size=14)

price_y = table_y + row_h * (len(rows) + 1)
add_cell(slide, MARGIN_X, price_y, feature_w, row_h, "PRICE", fill_color=BG, text_color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.LEFT, font_size=11)
price_values = [("FREE", ACCENT_GREEN), ("$4.99/mo", ACCENT_PURPLE), ("$2.99/pack", TEXT_PRIMARY), ("$9.99/mo", TEXT_PRIMARY)]
for c_idx, (value, color) in enumerate(price_values):
    x = MARGIN_X + feature_w + col_w * c_idx
    add_cell(slide, x, price_y, col_w, row_h, value, fill_color=BG, text_color=color, bold=True, align=PP_ALIGN.CENTER, font_size=12)

card_y = price_y + Inches(0.45)
card = add_box(slide, MARGIN_X, card_y, CONTENT_W, Inches(0.5), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), card_y + Inches(0.1), CONTENT_W - Inches(0.4), Inches(0.3), "Lifetime License: $49.99 one-time for permanent access to all premium features", font_name=FONT_BODY, font_size=12, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# Slide 6: Option B - 6 Buckets (with Travel)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_option_label(slide, "Slide 9 - Option B: 6 Buckets (with Travel)", MARGIN_X, Inches(0.35))
add_section_label(slide, "Funding Request", MARGIN_X, Inches(0.8))
add_h1(slide, "$4,500 to Validate, Polish, and Launch to Market", MARGIN_X, Inches(1.3), CONTENT_W)

wrap = add_box(slide, MARGIN_X, Inches(2.0), CONTENT_W, Inches(4.1), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(2.1), Inches(2.0), Inches(0.6), "$4,500", font_name=FONT_SANS, font_size=38, bold=True, color=ACCENT_GREEN)
add_textbox(slide, MARGIN_X + Inches(2.4), Inches(2.25), Inches(7.5), Inches(0.3), "Aligned with MIT Sandbox reimbursable categories", font_name=FONT_BODY, font_size=12, color=TEXT_SECONDARY)

items = [
    ("User Research", 1100, ACCENT_GREEN, "Run 40-60 customer discovery interviews with beta testers, including focus group sessions and small participation incentives."),
    ("Contractors", 1000, ACCENT_PURPLE, "Hire UI/UX designer to polish onboarding flow and create professional marketing assets for App Store launch."),
    ("Software & Tools", 800, ACCENT_PURPLE, "Subscribe to analytics tools (Mixpanel/Amplitude), design software (Figma), and cover Firebase scaling costs during beta."),
    ("Marketing & Launch", 600, ACCENT_PURPLE, "Launch targeted social media campaigns on Instagram/TikTok and create promotional video content for App Store listing."),
    ("Legal & Compliance", 500, ACCENT_PURPLE, "File for company incorporation, register as a legal entity, and secure registered agent services."),
    ("Travel & Conferences", 500, ACCENT_GREEN, "Attend wellness/productivity conferences to conduct customer research and connect with potential partners and influencers."),
]
bar_y = Inches(2.9)
row_gap = Inches(0.55)
bar_w = Inches(6.4)
for i, (label, amt, color, desc) in enumerate(items):
    y = bar_y + row_gap * i
    add_textbox(slide, MARGIN_X + Inches(0.2), y, Inches(1.8), Inches(0.3), label, font_name=FONT_BODY, font_size=10, bold=True, color=TEXT_PRIMARY)
    bar = add_box(slide, MARGIN_X + Inches(2.1), y + Inches(0.05), bar_w, Inches(0.2), border=TEXT_PRIMARY, fill_color=BG_SECONDARY, line_width=1)
    fill_w = bar_w * (amt / 4500)
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X + Inches(2.1), y + Inches(0.05), fill_w, Inches(0.2))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    add_textbox(slide, MARGIN_X + Inches(8.7), y, Inches(1.0), Inches(0.3), f"${amt}", font_name=FONT_BODY, font_size=10, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, MARGIN_X + Inches(0.2), y + Inches(0.23), CONTENT_W - Inches(0.4), Inches(0.25), desc, font_name=FONT_BODY, font_size=9, color=TEXT_SECONDARY)

note = add_box(slide, MARGIN_X, Inches(6.25), CONTENT_W, Inches(0.6), border=TEXT_PRIMARY, fill_color=BG, line_width=1)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(6.35), CONTENT_W - Inches(0.4), Inches(0.3), "Contractor spend within 25% cap. Legal within $1,000 cap. Conference fees within $500 cap. Travel requires pre-approved cost proposal per Sandbox guidelines.", font_name=FONT_BODY, font_size=10, color=TEXT_SECONDARY)

# Slide 8: Option B - Device Frames
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo(slide, SLIDE_W - Inches(1.2), Inches(0.35), Inches(0.7))
add_option_label(slide, "Demo Slide - Option B: Device Frames", MARGIN_X, Inches(0.35))
add_section_label(slide, "Product Demo", MARGIN_X, Inches(0.8))
add_h1(slide, "Experience Axiom Forge in Action", MARGIN_X, Inches(1.3), CONTENT_W)

card_y = Inches(2.3)
card_h = Inches(3.8)
for i, (title, border_color, play_color, label_text) in enumerate([
    ("Onboarding Flow", ACCENT_GREEN, ACCENT_GREEN, "Watch: 45 sec"),
    ("App Features", ACCENT_PURPLE, ACCENT_PURPLE, "Watch: 2 min"),
]):
    x = MARGIN_X + i * (COL_W + GAP_COL)
    add_h3(slide, title, x, card_y, COL_W, color=TEXT_PRIMARY)
    frame_w = Inches(3.4)
    frame_h = Inches(3.0)
    frame_x = x + (COL_W - frame_w) / 2
    frame_y = card_y + Inches(0.4)
    frame = add_box(slide, frame_x, frame_y, frame_w, frame_h, border=TEXT_PRIMARY, fill_color=BG, line_width=2)
    screen = add_phone_mockup(slide, frame_x + Inches(0.4), frame_y + Inches(0.2), Inches(2.6), Inches(2.6), border_color, play_color)
    add_textbox(slide, frame_x + Inches(0.4), frame_y + Inches(2.3), Inches(2.6), Inches(0.3), label_text, font_name=FONT_BODY, font_size=10, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

cta_y = Inches(6.1)
cta = add_box(slide, MARGIN_X, cta_y, CONTENT_W, Inches(0.9), border=TEXT_PRIMARY, fill_color=BG, line_width=2)
add_textbox(slide, MARGIN_X + Inches(0.3), cta_y + Inches(0.1), Inches(3.0), Inches(0.3), "Try It Yourself", font_name=FONT_SANS, font_size=16, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, MARGIN_X + Inches(0.3), cta_y + Inches(0.45), Inches(6.0), Inches(0.3), "No download required - works in your browser", font_name=FONT_BODY, font_size=11, color=TEXT_SECONDARY)
link = add_textbox(slide, MARGIN_X + Inches(8.8), cta_y + Inches(0.25), Inches(2.0), Inches(0.3), "Launch App", font_name=FONT_BODY, font_size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)
link.text_frame.paragraphs[0].runs[0].font.underline = True

prs.save("/workspaces/DMDHW/MIT_Sandbox/Pitch Deck/slide-options-v3.pptx")
print("Saved slide-options-v3 PPTX")
