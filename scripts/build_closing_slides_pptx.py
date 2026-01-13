from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.text import MSO_AUTO_SIZE

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0, 0, 0)
BG_SECONDARY = RGBColor(10, 10, 10)
BG_TERMINAL = RGBColor(10, 10, 10)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(192, 192, 192)
TEXT_MUTED = RGBColor(102, 102, 102)
ACCENT_PURPLE = RGBColor(167, 139, 250)
ACCENT_GREEN = RGBColor(204, 255, 0)
ACCENT_RED = RGBColor(255, 0, 0)

# Fonts
FONT_SANS = "Arial Black"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"

# Layout constants
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.55)
CONTENT_W = SLIDE_W - (MARGIN_X * 2)

def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_textbox(slide, x, y, w, h, text, font_name=FONT_BODY, font_size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box

def add_box(slide, x, y, w, h, border=None, fill_color=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
        
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_logo_shape(slide, x, y, size):
    # Draw a simplified version of the logo using shapes
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    base.fill.solid()
    base.fill.fore_color.rgb = TEXT_PRIMARY
    base.line.color.rgb = TEXT_PRIMARY
    base.line.width = Pt(1)
    # Bars
    bar_h = size * 0.16
    bar_w1 = size * 0.7
    bar_w2 = size * 0.45
    bar_w3 = size * 0.18
    offset_x = x + size * 0.15
    base_y = y + size * 0.62
    mid_y = y + size * 0.42
    top_y = y + size * 0.22
    bars = [
        (offset_x, base_y, bar_w1, bar_h),
        (x + size * 0.28, mid_y, bar_w2, bar_h),
        (x + size * 0.41, top_y, bar_w3, bar_h),
    ]
    for bx, by, bw, bh in bars:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bw, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = BG
        bar.line.fill.background()

def add_cta_button(slide, x, y, text, primary=True):
    w = Inches(3.0)
    h = Inches(0.8)
    if primary:
        btn = add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=TEXT_PRIMARY, line_width=2)
        text_color = BG
    else:
        btn = add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=None, line_width=2)
        text_color = TEXT_PRIMARY
        
    tf = btn.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_MONO
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_color
    return btn

def add_mindset_card(slide, x, y, w, h):
    # Card Background
    card = add_box(slide, x, y, w, h, border=TEXT_PRIMARY, fill_color=BG_SECONDARY, line_width=4)
    
    # Header Line
    header_y = y + Inches(0.8)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, header_y, w, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(50,50,50)
    line.line.fill.background()
    
    # Label
    add_textbox(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.3), "DAILY MENTALITY PILL", font_name=FONT_SANS, font_size=10, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, x + w - Inches(0.6), y + Inches(0.2), Inches(0.4), Inches(0.4), "💊", font_size=20)
    
    # Content
    content_y = header_y + Inches(0.4)
    # Quote Mark
    add_textbox(slide, x + Inches(0.3), content_y, Inches(0.5), Inches(0.5), "\"", font_name=FONT_SANS, font_size=60, bold=True, color=ACCENT_GREEN)
    
    # Quote Text
    quote = "Here's the big challenge of life. You can have more than you've got because you can become more than you are."
    add_textbox(slide, x + Inches(0.3), content_y + Inches(0.6), w - Inches(0.6), Inches(1.5), quote.upper(), font_name=FONT_MONO, font_size=16, bold=True, color=TEXT_PRIMARY)
    
    quote2 = "And of course the other side of the coin reads, unless you change how you are, you'll always have what you got."
    add_textbox(slide, x + Inches(0.3), content_y + Inches(2.0), w - Inches(0.6), Inches(1.0), quote2.upper(), font_name=FONT_MONO, font_size=12, bold=False, color=TEXT_SECONDARY)
    
    # Author Box
    auth_y = y + h - Inches(1.2)
    auth_box = add_box(slide, x + Inches(0.3), auth_y, w - Inches(0.6), Inches(0.9), fill_color=RGBColor(21,21,21))
    # Green strip
    strip = add_box(slide, x + Inches(0.3), auth_y, Inches(0.05), Inches(0.9), fill_color=ACCENT_GREEN)
    
    # Avatar square
    av_size = Inches(0.5)
    av = add_box(slide, x + Inches(0.5), auth_y + Inches(0.2), av_size, av_size, border=TEXT_PRIMARY, fill_color=ACCENT_GREEN, line_width=1)
    add_textbox(slide, x + Inches(0.5), auth_y + Inches(0.2), av_size, av_size, "JR", font_name=FONT_SANS, font_size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, x + Inches(1.2), auth_y + Inches(0.15), Inches(2), Inches(0.3), "JIM ROHN", font_name=FONT_SANS, font_size=12, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, x + Inches(1.2), auth_y + Inches(0.45), Inches(2), Inches(0.3), "ENTREPRENEUR & SPEAKER", font_name=FONT_BODY, font_size=9, color=TEXT_SECONDARY)


# ========================================== 
# OPTION 1: SPLIT SCREEN
# ========================================== 
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_shape(slide, SLIDE_W - Inches(1.0), Inches(0.5), Inches(0.8))
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3), "OPTION 1: SPLIT SCREEN", font_size=10, color=TEXT_MUTED)

# Left Side
add_textbox(slide, MARGIN_X, Inches(2.0), Inches(5), Inches(0.3), "// END OF PRESENTATION", font_name=FONT_MONO, font_size=12, bold=True, color=ACCENT_PURPLE)
add_textbox(slide, MARGIN_X, Inches(2.5), Inches(6), Inches(2), "THANK YOU FOR ARRIVING THIS FAR.", font_name=FONT_SANS, font_size=44, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, MARGIN_X, Inches(4.5), Inches(5.5), Inches(1), "Axiom Forge is ready to use. Experience the system that turns intent into identity.", font_name=FONT_BODY, font_size=16, color=TEXT_SECONDARY)
add_cta_button(slide, MARGIN_X, Inches(6.0), "TRY THE PROTOTYPE")

# Right Side
right_center = SLIDE_W * 0.75
add_textbox(slide, right_center - Inches(3), Inches(1.5), Inches(6), Inches(0.5), "// SPOILER ALERT: YOUR FIRST MINDSET PILL", font_name=FONT_MONO, font_size=10, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_mindset_card(slide, right_center - Inches(2.5), Inches(2.2), Inches(5), Inches(4.5))


# ========================================== 
# OPTION 2: PRODUCT CONTEXT (Phone)
# ========================================== 
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_shape(slide, SLIDE_W - Inches(1.0), Inches(0.5), Inches(0.8))
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3), "OPTION 2: PRODUCT CONTEXT", font_size=10, color=TEXT_MUTED)

# Left Side (Phone)
phone_w = Inches(3.5)
phone_h = Inches(6.5)
phone_x = Inches(2.0)
phone_y = Inches(0.8)

# Phone body
phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x, phone_y, phone_w, phone_h)
phone.fill.solid()
phone.fill.fore_color.rgb = BG
phone.line.color.rgb = RGBColor(50,50,50)
phone.line.width = Pt(4)

# Screen
screen_margin = Inches(0.15)
screen_x = phone_x + screen_margin
screen_y = phone_y + screen_margin
screen_w = phone_w - (screen_margin*2)
screen_h = phone_h - (screen_margin*2)
screen = add_box(slide, screen_x, screen_y, screen_w, screen_h, fill_color=BG)

# Phone Notch
notch_w = Inches(1.2)
notch_h = Inches(0.3)
notch = add_box(slide, phone_x + (phone_w - notch_w)/2, phone_y, notch_w, notch_h, fill_color=RGBColor(50,50,50))

# Phone Content
header_y = screen_y + Inches(0.5)
add_textbox(slide, screen_x, header_y, screen_w, Inches(0.3), "MINDSET", font_name=FONT_SANS, font_size=12, bold=True, align=PP_ALIGN.CENTER)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, screen_x, header_y + Inches(0.4), screen_w, Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(50,50,50)

# Phone - Pill
pill_y = header_y + Inches(0.6)
add_textbox(slide, screen_x + Inches(0.2), pill_y, screen_w, Inches(0.2), "// DAILY PILL", font_name=FONT_MONO, font_size=8, bold=True, color=ACCENT_GREEN)
add_textbox(slide, screen_x + Inches(0.2), pill_y + Inches(0.3), screen_w - Inches(0.4), Inches(1.5), "\"You can have more than you've got because you can become more than you are.\"", font_name=FONT_MONO, font_size=12, bold=True, color=TEXT_PRIMARY)

# Phone - Author
auth_y = pill_y + Inches(1.8)
auth_bg = add_box(slide, screen_x + Inches(0.2), auth_y, screen_w - Inches(0.4), Inches(0.5), fill_color=RGBColor(20,20,20))
strip = add_box(slide, screen_x + Inches(0.2), auth_y, Inches(0.05), Inches(0.5), fill_color=ACCENT_GREEN)
add_textbox(slide, screen_x + Inches(0.6), auth_y + Inches(0.1), Inches(2), Inches(0.3), "JIM ROHN", font_name=FONT_SANS, font_size=10, bold=True)
add_box(slide, screen_x + Inches(0.35), auth_y + Inches(0.12), Inches(0.25), Inches(0.25), fill_color=ACCENT_GREEN) # Mini avatar

# Phone - Reflection
ref_y = auth_y + Inches(0.8)
add_textbox(slide, screen_x + Inches(0.2), ref_y, screen_w, Inches(0.2), "// REFLECTION", font_name=FONT_MONO, font_size=8, bold=True, color=TEXT_MUTED)
add_box(slide, screen_x + Inches(0.2), ref_y + Inches(0.3), screen_w - Inches(0.4), Inches(0.6), border=RGBColor(50,50,50), fill_color=BG, line_width=1)
add_textbox(slide, screen_x + Inches(0.3), ref_y + Inches(0.4), screen_w, Inches(0.2), "How does this apply to you?", font_size=8, color=TEXT_MUTED)
add_textbox(slide, screen_x + Inches(0.2), ref_y + Inches(1.0), screen_w - Inches(0.4), Inches(0.3), "★ ★ ★ ★ ★", font_size=12, color=RGBColor(50,50,50), align=PP_ALIGN.CENTER)

# Right Side
right_x = SLIDE_W / 2 + Inches(0.5)
add_textbox(slide, right_x, Inches(2.5), Inches(5), Inches(1.5), "Start building your\nAxiom today.", font_name=FONT_SANS, font_size=40, bold=True)
# Color 'Axiom' green logic (simplified: overwrite text)
# Not easily doable in one run with simplified helper, but "Start building your..." is fine.

add_textbox(slide, right_x, Inches(4.2), Inches(5), Inches(1), "Join the closed beta to unlock the full system, including daily mindset pills and smart routine tracking.", font_name=FONT_BODY, font_size=16, color=TEXT_SECONDARY)
add_cta_button(slide, right_x, Inches(5.5), "LAUNCH PROTOTYPE")

# Footer Contact
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, Inches(6.5), Inches(5), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(50,50,50)
add_textbox(slide, right_x, Inches(6.6), Inches(5), Inches(0.8), "Javier Serrano\nMIT Sandbox\njavier@axiomforge.app", font_size=12, color=TEXT_SECONDARY)


# ========================================== 
# OPTION 3: PHILOSOPHY FIRST
# ========================================== 
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3), "OPTION 3: PHILOSOPHY FIRST", font_size=10, color=TEXT_MUTED)

center_x = SLIDE_W / 2

add_textbox(slide, MARGIN_X, Inches(1.5), CONTENT_W, Inches(0.5), "// YOUR FIRST PRINCIPLE", font_name=FONT_MONO, font_size=14, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)

# Quotes (Corners)
corner_size = Inches(0.8)
corner_thick = Pt(4)
# Top Left
tl_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(4), Inches(2.2), Pt(4), corner_size)
tl_v.fill.solid(); tl_v.fill.fore_color.rgb = ACCENT_GREEN
tl_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(4), Inches(2.2), corner_size, Pt(4))
tl_h.fill.solid(); tl_h.fill.fore_color.rgb = ACCENT_GREEN

# Bottom Right
br_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x + Inches(4), Inches(4.5) - corner_size, Pt(4), corner_size)
br_v.fill.solid(); br_v.fill.fore_color.rgb = ACCENT_GREEN
br_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x + Inches(4) - corner_size, Inches(4.5), corner_size, Pt(4))
br_h.fill.solid(); br_h.fill.fore_color.rgb = ACCENT_GREEN

# Main Text
add_textbox(slide, center_x - Inches(4), Inches(2.5), Inches(8), Inches(2), "\"UNLESS YOU CHANGE HOW YOU ARE, YOU'LL ALWAYS HAVE WHAT YOU GOT.\"", font_name=FONT_SANS, font_size=32, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, center_x - Inches(4), Inches(5.0), Inches(8), Inches(0.5), "—JIM ROHN", font_name=FONT_BODY, font_size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_textbox(slide, center_x - Inches(4), Inches(6.0), Inches(8), Inches(0.5), "Thank you for your time.", font_size=18, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

add_cta_button(slide, center_x - Inches(3.2), Inches(6.5), "TRY AXIOM FORGE")
btn2 = add_cta_button(slide, center_x + Inches(0.2), Inches(6.5), "CONTACT FOUNDER", primary=False)


# ========================================== 
# OPTION 4: SYSTEM TERMINAL
# ========================================== 
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_shape(slide, SLIDE_W - Inches(1.0), Inches(0.5), Inches(0.8))
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3), "OPTION 4: SYSTEM TERMINAL", font_size=10, color=TEXT_MUTED)

add_textbox(slide, MARGIN_X, Inches(1.0), CONTENT_W, Inches(0.8), "SYSTEM INITIALIZATION COMPLETE", font_name=FONT_SANS, font_size=36, bold=True, align=PP_ALIGN.CENTER)

term_w = Inches(9)
term_h = Inches(4.5)
term_x = (SLIDE_W - term_w) / 2
term_y = Inches(2.0)

# Window Body
add_box(slide, term_x, term_y, term_w, term_h, border=RGBColor(50,50,50), fill_color=BG_TERMINAL, line_width=2)
# Header
header_h = Inches(0.4)
add_box(slide, term_x, term_y, term_w, header_h, fill_color=RGBColor(50,50,50))
# Dots
dot_y = term_y + Inches(0.12)
for i, color in enumerate([RGBColor(255, 95, 86), RGBColor(255, 189, 46), RGBColor(39, 201, 63)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, term_x + Inches(0.2) + i*Inches(0.25), dot_y, Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()

add_textbox(slide, term_x + Inches(1.0), term_y + Inches(0.05), Inches(3), Inches(0.3), "axiom-forge -- init", font_name=FONT_MONO, font_size=10, color=TEXT_SECONDARY)

# Terminal Content
lines = [
    ("➜ loading routines... [OK]", TEXT_SECONDARY),
    ("➜ connecting behavioral_insights... [OK]", TEXT_SECONDARY),
    ("➜ generating mindset_pill_001... [DONE]", TEXT_SECONDARY),
    ("", TEXT_PRIMARY),
    ("OUTPUT RECEIVED:", ACCENT_GREEN),
]
current_y = term_y + header_h + Inches(0.2)
for text, color in lines:
    if text:
        add_textbox(slide, term_x + Inches(0.3), current_y, term_w, Inches(0.3), text, font_name=FONT_MONO, font_size=12, color=color)
    current_y += Inches(0.25)

# Output Box
out_x = term_x + Inches(0.3)
out_w = term_w - Inches(0.6)
add_box(slide, out_x, current_y, Inches(0.05), Inches(1.2), fill_color=TEXT_PRIMARY) # Left border
add_textbox(slide, out_x + Inches(0.2), current_y, out_w, Inches(0.8), "\"Here's the big challenge of life. You can have more than you've got because you can become more than you are.\"", font_name=FONT_MONO, font_size=12, color=TEXT_PRIMARY)
add_textbox(slide, out_x + Inches(0.2), current_y + Inches(0.8), out_w, Inches(0.3), ">> Jim Rohn, Entrepreneur", font_name=FONT_MONO, font_size=10, color=ACCENT_PURPLE)

current_y += Inches(1.5)
add_textbox(slide, term_x + Inches(0.3), current_y, Inches(2), Inches(0.3), "➜ ready for user input_", font_name=FONT_MONO, font_size=12, color=TEXT_SECONDARY)
# Cursor
add_box(slide, term_x + Inches(2.6), current_y, Inches(0.12), Inches(0.25), fill_color=ACCENT_GREEN)

add_cta_button(slide, (SLIDE_W - Inches(3))/2, term_y + term_h + Inches(0.3), "LAUNCH APPLICATION")


# ========================================== 
# OPTION 5: CARD REVEAL
# ========================================== 
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_logo_shape(slide, SLIDE_W - Inches(1.0), Inches(0.5), Inches(0.8))
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3), "OPTION 5: CARD REVEAL", font_size=10, color=TEXT_MUTED)

# Left Side
add_textbox(slide, MARGIN_X, Inches(2.0), Inches(5), Inches(0.3), "// PROTOTYPE READY", font_name=FONT_MONO, font_size=12, bold=True, color=ACCENT_GREEN)
add_textbox(slide, MARGIN_X, Inches(2.5), Inches(6), Inches(1.5), "FROM THEORY TO PRACTICE.", font_name=FONT_SANS, font_size=44, bold=True, color=TEXT_PRIMARY)
add_textbox(slide, MARGIN_X, Inches(4.2), Inches(5.5), Inches(1), "We've built the engine. Now we need the fuel.\nThank you for reviewing Axiom Forge.", font_name=FONT_BODY, font_size=16, color=TEXT_SECONDARY)

# Signature
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(5.5), Pt(2), Inches(0.8))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(50,50,50)
add_textbox(slide, MARGIN_X + Inches(0.2), Inches(5.5), Inches(4), Inches(0.8), "Javier Serrano\nMIT Sandbox | Winter 2026", font_size=12, color=TEXT_SECONDARY)

add_cta_button(slide, MARGIN_X, Inches(6.5), "OPEN PROTOTYPE")

# Right Side Stack
stack_center_x = SLIDE_W * 0.75
stack_center_y = Inches(3.75)
card_w = Inches(4.5)
card_h = Inches(2.8)

# Card 3 (Bottom)
c3 = add_box(slide, stack_center_x - card_w/2 + Inches(0.4), stack_center_y - card_h/2 + Inches(0.2), card_w, card_h, border=RGBColor(50,50,50), fill_color=RGBColor(20,20,20), line_width=2)
c3.rotation = -6

# Card 2 (Middle)
c2 = add_box(slide, stack_center_x - card_w/2 + Inches(0.2), stack_center_y - card_h/2 + Inches(0.1), card_w, card_h, border=RGBColor(80,80,80), fill_color=RGBColor(30,30,30), line_width=2)
c2.rotation = -3

# Card 1 (Top)
c1 = add_box(slide, stack_center_x - card_w/2, stack_center_y - card_h/2, card_w, card_h, border=TEXT_PRIMARY, fill_color=BG, line_width=4)
# Content inside Top Card
cx = stack_center_x - card_w/2
cy = stack_center_y - card_h/2
add_textbox(slide, cx + Inches(0.2), cy + Inches(0.1), Inches(1.5), Inches(0.3), "MINDSET #001", font_name=FONT_SANS, font_size=8, bold=True, color=BG).fill.solid(); 
# Bg for label
lbg = add_box(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(1.2), Inches(0.25), fill_color=ACCENT_GREEN)
add_textbox(slide, cx + Inches(0.25), cy + Inches(0.15), Inches(1.2), Inches(0.25), "MINDSET #001", font_name=FONT_SANS, font_size=9, bold=True, color=BG)

add_textbox(slide, cx + Inches(3.8), cy + Inches(0.15), Inches(0.5), Inches(0.3), "•••", color=TEXT_SECONDARY)

add_textbox(slide, cx + Inches(0.3), cy + Inches(0.6), card_w - Inches(0.6), Inches(1.5), "\"YOU CAN HAVE MORE THAN YOU'VE GOT BECAUSE YOU CAN BECOME MORE THAN YOU ARE.\"", font_name=FONT_MONO, font_size=14, bold=True, color=TEXT_PRIMARY)

# Footer inside card
cf_y = cy + card_h - Inches(0.6)
add_box(slide, cx + Inches(0.3), cf_y, Inches(0.3), Inches(0.3), border=TEXT_PRIMARY, line_width=1)
add_textbox(slide, cx + Inches(0.3), cf_y, Inches(0.3), Inches(0.3), "JR", font_size=8, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, cx + Inches(0.7), cf_y, Inches(2), Inches(0.3), "JIM ROHN", font_size=10, color=TEXT_SECONDARY)


# Save
output_path = "/workspaces/DMDHW/MIT_Sandbox/Pitch Deck/closing-slide-options.pptx"
prs.save(output_path)
print(f"Saved PPTX to {output_path}")
